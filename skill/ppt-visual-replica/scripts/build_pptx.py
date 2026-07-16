#!/usr/bin/env python
from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ALLOWED_NATIVE = {
    "background",
    "panel",
    "frame",
    "divider",
    "connector",
    "structural_arrow",
    "text_container",
    "layout_only",
}
ALLOWED_IMAGE_SOURCES = {"imagegen_asset", "openai_image", "gemini_image", "user_asset"}
ARROWHEAD_TYPES = {"none", "triangle", "stealth", "diamond", "oval", "arrow"}
ARROWHEAD_SIZES = {"sm", "med", "lg"}


def require(condition: bool, message: str) -> None:
    if not condition:
        raise ValueError(message)


def rgb(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    require(len(value) == 6, f"invalid RGB color: {value}")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def emu(canvas_units: float, scale: float) -> int:
    return int(Inches(canvas_units * scale))


def element_name(item: dict[str, Any]) -> str:
    value = item.get("id")
    require(bool(value), "every layout element must declare id")
    return str(value)


def set_shape_name(shape: Any, name: str) -> None:
    candidates = (
        getattr(getattr(shape._element, "nvPicPr", None), "cNvPr", None),
        getattr(getattr(shape._element, "nvSpPr", None), "cNvPr", None),
        getattr(getattr(shape._element, "nvCxnSpPr", None), "cNvPr", None),
    )
    node = next((candidate for candidate in candidates if candidate is not None), None)
    require(node is not None, f"cannot locate non-visual properties for {name}")
    node.set("name", name)


def set_picture_metadata(shape: Any, item: dict[str, Any]) -> None:
    semantic_unit_id = str(item["semantic_unit_id"])
    metadata = {
        "semantic_unit_id": semantic_unit_id,
        "semantic_unit_count": int(item["semantic_unit_count"]),
        "source_type": str(item["source_type"]),
        "asset_role": item.get("asset_role", "semantic_unit"),
    }
    node = shape._element.nvPicPr.cNvPr
    node.set("name", f"semantic:{semantic_unit_id}")
    node.set("descr", json.dumps(metadata, ensure_ascii=False, separators=(",", ":")))


def arrowhead_spec(item: dict[str, Any], key: str) -> dict[str, str] | None:
    raw = item.get(key)
    if raw in (None, False, "", "none"):
        return None
    if isinstance(raw, dict):
        arrow_type = str(raw.get("type") or item.get("arrow_type") or "triangle")
        width = str(raw.get("width") or raw.get("w") or item.get("arrow_width") or "sm")
        length = str(
            raw.get("length") or raw.get("len") or item.get("arrow_length") or "sm"
        )
    else:
        arrow_type = str(
            raw if isinstance(raw, str) else item.get("arrow_type") or "triangle"
        )
        width = str(item.get("arrow_width") or "sm")
        length = str(item.get("arrow_length") or "sm")
    require(
        arrow_type in ARROWHEAD_TYPES,
        f"{item.get('id')}: unsupported arrowhead type: {arrow_type}",
    )
    require(
        width in ARROWHEAD_SIZES,
        f"{item.get('id')}: unsupported arrowhead width: {width}",
    )
    require(
        length in ARROWHEAD_SIZES,
        f"{item.get('id')}: unsupported arrowhead length: {length}",
    )
    return {"type": arrow_type, "w": width, "len": length}


def set_native_arrowheads(shape: Any, item: dict[str, Any]) -> None:
    line = shape._element.spPr.get_or_add_ln()
    for tag in ("a:headEnd", "a:tailEnd"):
        existing = line.find(qn(tag))
        if existing is not None:
            line.remove(existing)
    for key, tag in (("begin_arrow", "a:headEnd"), ("end_arrow", "a:tailEnd")):
        spec = arrowhead_spec(item, key)
        if not spec:
            continue
        node = OxmlElement(tag)
        node.set("type", spec["type"])
        node.set("w", spec["w"])
        node.set("len", spec["len"])
        line.append(node)


def scene_exception_is_valid(item: dict[str, Any]) -> bool:
    if (
        item.get("asset_role") != "generated_scene_background"
        or item.get("allows_text_overlay") is not True
    ):
        return False
    exception = item.get("composite_exception")
    if exception is None:
        require(
            item.get("contains_separable_foreground") is False,
            f"{item.get('id')}: scene background must declare contains_separable_foreground: false or a user-approved composite_exception",
        )
        return True
    require(
        isinstance(exception, dict),
        f"{item.get('id')}: composite_exception must be an object",
    )
    require(
        str(exception.get("approved_by")).lower() == "user",
        f"{item.get('id')}: composite exception must be user-approved",
    )
    require(
        bool(exception.get("scope")) and bool(exception.get("reason")),
        f"{item.get('id')}: composite exception needs scope and reason",
    )
    return True


def image_box(
    item: dict[str, Any], asset_dir: Path
) -> tuple[float, float, float, float]:
    slot = item.get("anchor_slot") or [
        item.get("x"),
        item.get("y"),
        item.get("w"),
        item.get("h"),
    ]
    require(
        isinstance(slot, list) and len(slot) == 4,
        f"{item.get('id')}: image anchor slot must contain four values",
    )
    x, y, width, height = [float(value) for value in slot]
    require(
        width > 0 and height > 0,
        f"{item.get('id')}: image anchor slot must be positive",
    )
    path = asset_dir / str(item["path"])
    require(path.is_file(), f"{item.get('id')}: image asset not found: {path}")
    with Image.open(path) as image:
        image_width, image_height = image.size
    require(
        image_width > 0 and image_height > 0, f"{item.get('id')}: image asset is empty"
    )
    fit = min(width / image_width, height / image_height)
    fitted_width, fitted_height = image_width * fit, image_height * fit
    return (
        x + (width - fitted_width) / 2,
        y + (height - fitted_height) / 2,
        fitted_width,
        fitted_height,
    )


def add_text(slide: Any, item: dict[str, Any], scale: float) -> None:
    name = element_name(item)
    text = str(item.get("text", ""))
    lines = text.splitlines() or [""]
    expected_lines = item.get("expected_lines")
    if expected_lines is not None:
        require(
            len(lines) == int(expected_lines),
            f"{name}: text line count does not match expected_lines",
        )
    shape = slide.shapes.add_textbox(
        emu(float(item["x"]), scale),
        emu(float(item["y"]), scale),
        emu(float(item["w"]), scale),
        emu(float(item["h"]), scale),
    )
    set_shape_name(shape, name)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(str(item.get("valign", "top")), MSO_ANCHOR.TOP)
    for margin_name in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(frame, margin_name, Pt(float(item.get(margin_name, 1))))
    for index, line_text in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }.get(str(item.get("align", "left")), PP_ALIGN.LEFT)
        run = paragraph.add_run()
        run.text = line_text
        run.font.name = str(item.get("font", "Microsoft YaHei"))
        run.font.size = Pt(float(item.get("font_size", 12)))
        run.font.bold = bool(item.get("bold", False))
        run.font.italic = bool(item.get("italic", False))
        run.font.color.rgb = rgb(str(item.get("color", "#111111")))
    if item.get("rotation") is not None:
        shape.rotation = float(item["rotation"])


def add_shape(slide: Any, item: dict[str, Any], scale: float) -> None:
    name = element_name(item)
    purpose = str(item.get("purpose", ""))
    require(
        purpose in ALLOWED_NATIVE,
        f"{name}: native shape purpose is not structural: {purpose}",
    )
    shape_type = {
        "rect": MSO_SHAPE.RECTANGLE,
        "round_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "right_arrow": MSO_SHAPE.RIGHT_ARROW,
        "down_arrow": MSO_SHAPE.DOWN_ARROW,
    }.get(str(item["type"]))
    require(shape_type is not None, f"{name}: unsupported native shape type")
    if item["type"] in {"right_arrow", "down_arrow"}:
        require(
            purpose == "structural_arrow",
            f"{name}: block arrow must have structural_arrow purpose",
        )
        require(
            item.get("block_arrow_reference") is True,
            f"{name}: block arrow requires block_arrow_reference: true",
        )
    shape = slide.shapes.add_shape(
        shape_type,
        emu(float(item["x"]), scale),
        emu(float(item["y"]), scale),
        emu(float(item["w"]), scale),
        emu(float(item["h"]), scale),
    )
    set_shape_name(shape, name)
    if item["type"] == "round_rect" and "radius" in item:
        try:
            shape.adjustments[0] = float(item["radius"])
        except (IndexError, ValueError):
            pass
    try:
        shape.shadow.inherit = False
    except (AttributeError, ValueError):
        pass
    if item.get("fill"):
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(str(item["fill"]))
    else:
        shape.fill.background()
    if item.get("line"):
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = rgb(str(item["line"]))
        shape.line.width = Pt(float(item.get("line_width", 1)))
    else:
        shape.line.fill.background()


def add_line(slide: Any, item: dict[str, Any], scale: float) -> None:
    name = element_name(item)
    purpose = str(item.get("purpose", ""))
    require(
        purpose in {"connector", "structural_arrow", "divider", "layout_only"},
        f"{name}: line purpose must be structural",
    )
    if purpose == "structural_arrow":
        require(
            item.get("begin_arrow") or item.get("end_arrow"),
            f"{name}: structural arrow must use native arrowhead metadata",
        )
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        emu(float(item["x1"]), scale),
        emu(float(item["y1"]), scale),
        emu(float(item["x2"]), scale),
        emu(float(item["y2"]), scale),
    )
    set_shape_name(shape, name)
    shape.line.fill.solid()
    shape.line.fill.fore_color.rgb = rgb(str(item.get("color", "#000000")))
    shape.line.width = Pt(float(item.get("width", 1)))
    dash = str(item.get("dash", "solid"))
    if dash != "solid":
        dash_map = {
            "dash": MSO_LINE_DASH_STYLE.DASH,
            "long_dash": MSO_LINE_DASH_STYLE.LONG_DASH,
            "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
            "square_dot": MSO_LINE_DASH_STYLE.SQUARE_DOT,
            "dash_dot": MSO_LINE_DASH_STYLE.DASH_DOT,
        }
        require(dash in dash_map, f"{name}: unsupported dash style: {dash}")
        shape.line.dash_style = dash_map[dash]
    set_native_arrowheads(shape, item)


def add_image(
    slide: Any, item: dict[str, Any], asset_dir: Path, scale: float
) -> list[float]:
    name = element_name(item)
    source_type = str(item.get("source_type", ""))
    require(
        source_type in ALLOWED_IMAGE_SOURCES,
        f"{name}: invalid image source_type: {source_type}",
    )
    scene = scene_exception_is_valid(item)
    if scene:
        semantic_unit_id = str(item.get("semantic_unit_id") or name)
        item["semantic_unit_id"] = semantic_unit_id
        item["semantic_unit_count"] = int(item.get("semantic_unit_count", 1))
    else:
        require(
            bool(item.get("semantic_unit_id")),
            f"{name}: image element must declare semantic_unit_id",
        )
        require(
            "semantic_unit_count" in item,
            f"{name}: image element must declare semantic_unit_count",
        )
        require(
            int(item["semantic_unit_count"]) == 1,
            f"{name}: final image asset must contain one minimum semantic unit",
        )
    x, y, width, height = image_box(item, asset_dir)
    picture = slide.shapes.add_picture(
        str(asset_dir / str(item["path"])),
        emu(x, scale),
        emu(y, scale),
        emu(width, scale),
        emu(height, scale),
    )
    set_picture_metadata(picture, item)
    fitted = [round(value, 3) for value in (x, y, width, height)]
    item["fitted_bbox"] = fitted
    item["scaling"] = "uniform_contain"
    return fitted


def build(manifest_path: Path, output_path: Path, asset_dir: Path) -> dict[str, Any]:
    manifest = json.loads(manifest_path.read_text(encoding="utf-8-sig"))
    canvas = manifest.get("canvas", {"width": 1600, "height": 900})
    slide_size = manifest.get("slide", {"width_in": 16, "height_in": 9})
    canvas_width = float(canvas["width"])
    canvas_height = float(canvas["height"])
    slide_width = float(slide_size["width_in"])
    slide_height = float(slide_size["height_in"])
    require(
        canvas_width > 0 and canvas_height > 0 and slide_width > 0 and slide_height > 0,
        "canvas and slide sizes must be positive",
    )
    scale_x = slide_width / canvas_width
    scale_y = slide_height / canvas_height
    require(
        abs(scale_x - scale_y) <= max(scale_x, scale_y) * 0.002,
        "canvas and slide aspect ratios do not match",
    )
    scale = (scale_x + scale_y) / 2

    presentation = Presentation()
    presentation.slide_width = Inches(slide_width)
    presentation.slide_height = Inches(slide_height)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    seen_ids: set[str] = set()
    counts = {"text": 0, "image": 0, "line": 0, "shape": 0}
    for item in manifest.get("elements", []):
        name = element_name(item)
        require(name not in seen_ids, f"duplicate layout element id: {name}")
        seen_ids.add(name)
        element_type = str(item.get("type", ""))
        if element_type == "text":
            add_text(slide, item, scale)
            counts["text"] += 1
        elif element_type == "image":
            add_image(slide, item, asset_dir, scale)
            counts["image"] += 1
        elif element_type == "line":
            add_line(slide, item, scale)
            counts["line"] += 1
        elif element_type in {"rect", "round_rect", "right_arrow", "down_arrow"}:
            add_shape(slide, item, scale)
            counts["shape"] += 1
        else:
            raise ValueError(f"{name}: unsupported element type: {element_type}")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    manifest["build"] = {
        "status": "passed",
        "pptx": str(output_path),
        "object_counts": counts,
        "image_scaling": "uniform_contain",
        "semantic_metadata_embedded": True,
    }
    return manifest


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Build a strict image-first PPTX replica."
    )
    parser.add_argument("--manifest", required=True, type=Path)
    parser.add_argument("--asset-dir", required=True, type=Path)
    parser.add_argument("--out", required=True, type=Path)
    parser.add_argument("--resolved-manifest-out", type=Path)
    args = parser.parse_args()

    resolved = build(args.manifest, args.out, args.asset_dir)
    resolved_path = args.resolved_manifest_out or args.manifest
    resolved_path.parent.mkdir(parents=True, exist_ok=True)
    resolved_path.write_text(
        json.dumps(resolved, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    print(
        json.dumps(
            {
                "status": "passed",
                "pptx": str(args.out),
                "resolved_manifest": str(resolved_path),
                "object_counts": resolved["build"]["object_counts"],
            },
            ensure_ascii=False,
        )
    )


if __name__ == "__main__":
    main()
